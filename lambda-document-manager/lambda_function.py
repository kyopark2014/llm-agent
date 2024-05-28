import json
import boto3
import os
import traceback
import PyPDF2
import time
import docx
import base64

from io import BytesIO
from urllib import parse
from botocore.config import Config
from PIL import Image
from urllib.parse import unquote_plus
from langchain_community.embeddings import BedrockEmbeddings
from langchain_community.vectorstores.opensearch_vector_search import OpenSearchVectorSearch
from langchain_community.docstore.document import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from opensearchpy import OpenSearch
from pptx import Presentation
from multiprocessing import Process, Pipe
from langchain_core.prompts import ChatPromptTemplate
from langchain_community.chat_models import BedrockChat
from langchain.callbacks.streaming_stdout import StreamingStdOutCallbackHandler
from langchain_core.messages import HumanMessage, SystemMessage
from langchain_aws import ChatBedrock

s3 = boto3.client('s3')

s3_bucket = os.environ.get('s3_bucket') # bucket name
s3_prefix = os.environ.get('s3_prefix')
meta_prefix = "metadata/"
profile_of_LLMs = json.loads(os.environ.get('profile_of_LLMs'))
enableParallelSummay = os.environ.get('enableParallelSummay')

opensearch_account = os.environ.get('opensearch_account')
opensearch_passwd = os.environ.get('opensearch_passwd')
opensearch_url = os.environ.get('opensearch_url')
sqsUrl = os.environ.get('sqsUrl')
doc_prefix = s3_prefix+'/'

sqs = boto3.client('sqs')
s3_client = boto3.client('s3')  

roleArn = os.environ.get('roleArn') 
path = os.environ.get('path')
max_object_size = int(os.environ.get('max_object_size'))

supportedFormat = json.loads(os.environ.get('supportedFormat'))
print('supportedFormat: ', supportedFormat)

os_client = OpenSearch(
    hosts = [{
        'host': opensearch_url.replace("https://", ""), 
        'port': 443
    }],
    http_compress = True,
    http_auth=(opensearch_account, opensearch_passwd),
    use_ssl = True,
    verify_certs = True,
    ssl_assert_hostname = False,
    ssl_show_warn = False,
)

def delete_document_if_exist(metadata_key):
    try: 
        s3r = boto3.resource("s3")
        bucket = s3r.Bucket(s3_bucket)
        objs = list(bucket.objects.filter(Prefix=metadata_key))
        print('objs: ', objs)
        
        if(len(objs)>0):
            doc = s3r.Object(s3_bucket, metadata_key)
            meta = doc.get()['Body'].read().decode('utf-8')
            print('meta: ', meta)
            
            ids = json.loads(meta)['ids']
            print('ids: ', ids)
            
            result = vectorstore.delete(ids)
            print('result: ', result)        
        else:
            print('no meta file: ', metadata_key)
            
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        raise Exception ("Not able to create meta file")

# embedding for RAG
region_name = os.environ.get('bedrock_region')
    
boto3_bedrock = boto3.client(
    service_name='bedrock-runtime',
    region_name=region_name,
    config=Config(
        retries = {
            'max_attempts': 30
        }            
    )
)

bedrock_embeddings = BedrockEmbeddings(
    client=boto3_bedrock,
    region_name = region_name,
    model_id = 'amazon.titan-embed-text-v1' 
)   

index_name = 'idx-rag'
vectorstore = OpenSearchVectorSearch(
    index_name=index_name,  
    is_aoss = False,
    #engine="faiss",  # default: nmslib
    embedding_function = bedrock_embeddings,
    opensearch_url = opensearch_url,
    http_auth=(opensearch_account, opensearch_passwd),
)  

def store_document_for_opensearch(docs, key):    
    objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
    print('objectName: ', objectName)    
    metadata_key = meta_prefix+objectName+'.metadata.json'
    print('meta file name: ', metadata_key)    
    delete_document_if_exist(metadata_key)
    
    try:        
        response = vectorstore.add_documents(docs, bulk_size = 2000)
        print('response of adding documents: ', response)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                
        #raise Exception ("Not able to request to LLM")

    print('uploaded into opensearch')
    
    return response
    
def delete_document_if_exist(metadata_key):
    try: 
        s3r = boto3.resource("s3")
        bucket = s3r.Bucket(s3_bucket)
        objs = list(bucket.objects.filter(Prefix=metadata_key))
        print('objs: ', objs)
        
        if(len(objs)>0):
            doc = s3r.Object(s3_bucket, metadata_key)
            meta = doc.get()['Body'].read().decode('utf-8')
            print('meta: ', meta)
            
            ids = json.loads(meta)['ids']
            print('ids: ', ids)
            
            result = vectorstore.delete(ids)
            print('result: ', result)        
        else:
            print('no meta file: ', metadata_key)
            
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        raise Exception ("Not able to create meta file")
             
# load documents from s3 for pdf and txt
def load_document(file_type, key):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, key)
    
    contents = ""
    if file_type == 'pdf':
        Byte_contents = doc.get()['Body'].read()
        
        try: 
            reader = PyPDF2.PdfReader(BytesIO(Byte_contents))
            
            texts = []
            for page in reader.pages:
                texts.append(page.extract_text())
            contents = '\n'.join(texts)
        except Exception:
                err_msg = traceback.format_exc()
                print('err_msg: ', err_msg)
                # raise Exception ("Not able to load the pdf file")
                
    elif file_type == 'pptx':
        Byte_contents = doc.get()['Body'].read()
            
        try:
            prs = Presentation(BytesIO(Byte_contents))

            texts = []
            for i, slide in enumerate(prs.slides):
                text = ""
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = text + shape.text
                texts.append(text)
            contents = '\n'.join(texts)
        except Exception:
                err_msg = traceback.format_exc()
                print('err_msg: ', err_msg)
                # raise Exception ("Not able to load texts from preseation file")
        
    elif file_type == 'txt' or file_type == 'md':       
        try:  
            contents = doc.get()['Body'].read().decode('utf-8')
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)        
            # raise Exception ("Not able to load the file")

    elif file_type == 'docx':
        try:
            Byte_contents = doc.get()['Body'].read()                    
            doc_contents =docx.Document(BytesIO(Byte_contents))

            texts = []
            for i, para in enumerate(doc_contents.paragraphs):
                if(para.text):
                    texts.append(para.text)
                    # print(f"{i}: {para.text}")        
            contents = '\n'.join(texts)            
            # print('contents: ', contents)
        except Exception:
                err_msg = traceback.format_exc()
                print('err_msg: ', err_msg)
                # raise Exception ("Not able to load docx")   
    
    texts = ""
    if len(contents)>0:
        new_contents = str(contents).replace("\n"," ") 
        print('length: ', len(new_contents))
        
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=100,
            separators=["\n\n", "\n", ".", " ", ""],
            length_function = len,
        ) 

        texts = text_splitter.split_text(new_contents) 
                        
    return texts

# load a code file from s3
def load_code(file_type, key):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, key)
    
    if file_type == 'py':        
        contents = doc.get()['Body'].read().decode('utf-8')
        separators = ["\ndef "]
        #print('contents: ', contents)
    elif file_type == 'js':
        contents = doc.get()['Body'].read().decode('utf-8')
        separators = ["\nfunction ", "\nexports.handler "]
    
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=50,
        chunk_overlap=0,
        #separators=["def ", "\n\n", "\n", ".", " ", ""],
        separators=separators,
        length_function = len,
    ) 

    texts = text_splitter.split_text(contents) 
    
    for i, text in enumerate(texts):
        print(f"Chunk #{i}: {text}")
                
    return texts

def isSupported(type):
    for format in supportedFormat:
        if type == format:
            return True
    
    return False
    
def check_supported_type(key, file_type, size):    
    if key.find('/html/') != -1 or key.find('/node_modules/') != -1 or key.find('/.git/') != -1: # do not include html/node_modules folders
        print('html: ', key.find('/html/'))
        return False
    
    if isSupported(file_type):
        if key[0]=='.' or key[key.rfind('/')+1]=='.':
            print(f"Ignore {key} since the filename starts a dot character for macbook.")        
            return False
        elif size > 0 and size<max_object_size:
            return True
    else:
        return False

HUMAN_PROMPT = "\n\nHuman:"
AI_PROMPT = "\n\nAssistant:"
def get_parameter(model_type, maxOutputTokens):
    if model_type=='titan': 
        return {
            "maxTokenCount":1024,
            "stopSequences":[],
            "temperature":0,
            "topP":0.9
        }
    elif model_type=='claude':
        return {
            "max_tokens_to_sample":maxOutputTokens, # 8k    
            "temperature":0.1,
            "top_k":250,
            "top_p":0.9,
            "stop_sequences": [HUMAN_PROMPT]            
        }
        
# Multi-LLM
def get_chat(profile_of_LLMs, selected_LLM):
    profile = profile_of_LLMs[selected_LLM]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'LLM: {selected_LLM}, bedrock_region: {bedrock_region}, modelId: {modelId}')
    maxOutputTokens = int(profile['maxOutputTokens'])
                          
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }            
        )
    )
    parameters = {
        "max_tokens":maxOutputTokens,     
        "temperature":0.1,
        "top_k":250,
        "top_p":0.9,
        "stop_sequences": [HUMAN_PROMPT]
    }
    # print('parameters: ', parameters)

    chat = ChatBedrock(  
        model_id=modelId,
        client=boto3_bedrock, 
        model_kwargs=parameters,
    )         
    
    return chat

def summary_of_code(chat, code, mode):
    if mode == 'py': 
        system = (
            "다음의 <article> tag에는 python code가 있습니다. code의 전반적인 목적에 대해 설명하고, 각 함수의 기능과 역할을 자세하게 한국어 500자 이내로 설명하세요."
        )
    elif mode == 'js':
        system = (
            "다음의 <article> tag에는 node.js code가 있습니다. code의 전반적인 목적에 대해 설명하고, 각 함수의 기능과 역할을 자세하게 한국어 500자 이내로 설명하세요."
        )
    else:
        system = (
            "다음의 <article> tag에는 code가 있습니다. code의 전반적인 목적에 대해 설명하고, 각 함수의 기능과 역할을 자세하게 한국어 500자 이내로 설명하세요."
        )
    
    human = "<article>{code}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "code": code
            }
        )
        
        summary = result.content
        print('result of code summarization: ', summary)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return summary

def summarize_process_for_relevent_code(conn, chat, code, key, region_name):
    try: 
        if code.find('\ndef ') != -1:
            start = code.find('\ndef ')
            end = code.find(':')   
        elif code.find('\nfunction ') != -1:
            start = code.find('\nfunction ')
            end = code.find('(')   
        elif code.find('\nexports.') != -1:
            start = code.find('\nexports.')
            end = code.find(' =')         
        else:
            start = -1
            end = -1
              
        print('code: ', code)                             
        print(f'start: {start}, end: {end}')
                    
        doc = ""    
        if start != -1:      
            function_name = code[start+1:end]
            print('function_name: ', function_name)
            
            file_type = key[key.rfind('.')+1:len(key)].lower()
            print('file_type: ', file_type)
                            
            summary = summary_of_code(chat, code, file_type)
            print(f"summary ({region_name}, {file_type}): {summary}")
            
            # print('first line summary: ', summary[:len(function_name)])
            # print('function name: ', function_name)            
            if summary[:len(function_name)]==function_name:
                summary = summary[summary.find('\n')+1:len(summary)]

            doc = Document(
                page_content=summary,
                metadata={
                    'name': key,
                    # 'uri': path+doc_prefix+parse.quote(key),
                    'uri': path+key,
                    'code': code,
                    'function_name': function_name
                }
            )           
                        
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)       
        # raise Exception (f"Not able to summarize: {doc}")               
    
    conn.send(doc)    
    conn.close()

def summarize_relevant_codes_using_parallel_processing(codes, key):
    selected_LLM = 0
    relevant_codes = []    
    processes = []
    parent_connections = []
    for code in codes:
        parent_conn, child_conn = Pipe()
        parent_connections.append(parent_conn)
            
        chat = get_chat(profile_of_LLMs, selected_LLM)
        region_name = profile_of_LLMs[selected_LLM]['bedrock_region']

        process = Process(target=summarize_process_for_relevent_code, args=(child_conn, chat, code, key, region_name))
        processes.append(process)

        selected_LLM = selected_LLM + 1
        if selected_LLM == len(profile_of_LLMs):
            selected_LLM = 0

    for process in processes:
        process.start()
            
    for parent_conn in parent_connections:
        doc = parent_conn.recv()
        
        if doc:
            relevant_codes.append(doc)    

    for process in processes:
        process.join()
    
    return relevant_codes

def extract_text(chat, img_base64):    
    query = "텍스트를 추출해서 utf8로 변환하세요. <result> tag를 붙여주세요."
    
    messages = [
        HumanMessage(
            content=[
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/png;base64,{img_base64}", 
                    },
                },
                {
                    "type": "text", "text": query
                },
            ]
        )
    ]
    
    try: 
        result = chat.invoke(messages)
        
        extracted_text = result.content
        print('result of text extraction from an image: ', extracted_text)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return extracted_text

def get_documentId(key, category):
    documentId = category + "-" + key
    documentId = documentId.replace(' ', '_') # remove spaces  
    documentId = documentId.replace(',', '_') # remove commas # not allowed: [ " * \\ < | , > / ? ]
    documentId = documentId.replace('/', '_') # remove slash
    documentId = documentId.lower() # change to lowercase
                
    return documentId

def create_metadata(bucket, key, meta_prefix, s3_prefix, uri, category, documentId, ids):
    title = key
    timestamp = int(time.time())

    metadata = {
        "Attributes": {
            "_category": category,
            "_source_uri": uri,
            "_version": str(timestamp),
            "_language_code": "ko"
        },
        "Title": title,
        "DocumentId": documentId,      
        "ids": ids  
    }
    print('metadata: ', metadata)
    
    #objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)]).upper()
    objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
    print('objectName: ', objectName)

    client = boto3.client('s3')
    try: 
        client.put_object(
            Body=json.dumps(metadata), 
            Bucket=bucket, 
            Key=meta_prefix+objectName+'.metadata.json' 
        )
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        raise Exception ("Not able to create meta file")
    
# load csv documents from s3
def lambda_handler(event, context):
    print('event: ', event)    
    
    documentIds = []
    for record in event['Records']:
        receiptHandle = record['receiptHandle']
        print("receiptHandle: ", receiptHandle)
        
        body = record['body']
        print("body: ", body)
        
        jsonbody = json.loads(body)        
        bucket = jsonbody['bucket']        
        # translate utf8
        key = unquote_plus(jsonbody['key']) # url decoding
        print('bucket: ', bucket)
        print('key: ', key)        
        eventName = jsonbody['type']
        
        start_time = time.time()      
        
        file_type = key[key.rfind('.')+1:len(key)].lower()
        print('file_type: ', file_type)
            
        if eventName == 'ObjectRemoved:Delete':
            if isSupported(file_type):
                objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
                print('objectName: ', objectName)
                
                # get metadata from s3
                metadata_key = meta_prefix+objectName+'.metadata.json'
                print('metadata_key: ', metadata_key)

                documentId = ""
                try: 
                    metadata_obj = s3.get_object(Bucket=bucket, Key=metadata_key)
                    metadata_body = metadata_obj['Body'].read().decode('utf-8')
                    metadata = json.loads(metadata_body)
                    print('metadata: ', metadata)
                    documentId = metadata['DocumentId']
                    print('documentId: ', documentId)
                    documentIds.append(documentId)
                except Exception:
                    err_msg = traceback.format_exc()
                    print('err_msg: ', err_msg)
                    # raise Exception ("Not able to get the object")
                    
                if documentId:
                    try: # delete metadata                        
                        delete_document_if_exist(metadata_key)
                        
                        print('delete metadata: ', metadata_key)                        
                        result = s3.delete_object(Bucket=bucket, Key=metadata_key)
                        # print('result of metadata deletion: ', result)
                        
                    except Exception:
                        err_msg = traceback.format_exc()
                        print('err_msg: ', err_msg)
                        # raise Exception ("Not able to delete documents in Kendra")                    
            else: 
                print('This file format is not supported: ', file_type)                
                    
        elif eventName == "ObjectCreated:Put":            
            size = 0
            try:
                s3obj = s3.get_object(Bucket=bucket, Key=key)
                print(f"Got object: {s3obj}")        
                size = int(s3obj['ContentLength'])    
                
                #attributes = ['ETag', 'Checksum', 'ObjectParts', 'StorageClass', 'ObjectSize']
                #result = s3.get_object_attributes(Bucket=bucket, Key=key, ObjectAttributes=attributes)  
                #print('result: ', result)            
                #size = int(result['ObjectSize'])
                print('object size: ', size)
            except Exception:
                err_msg = traceback.format_exc()
                print('err_msg: ', err_msg)
                # raise Exception ("Not able to get object info") 
            
            if check_supported_type(key, file_type, size): 
                if file_type == 'py' or file_type == 'js':  # for code
                    category = file_type
                #elif file_type == 'png' or file_type == 'jpg' or file_type == 'jpeg':
                #    category = 'img'
                else:
                    category = "upload" # for document
                documentId = get_documentId(key, category)                                
                print('documentId: ', documentId)
                
                docs = []
                        
                if file_type == 'pdf' or file_type == 'txt' or file_type == 'md' or file_type == 'csv' or file_type == 'pptx' or file_type == 'docx':
                    print('upload to opensearch: ', key) 
                    texts = load_document(file_type, key)
                            
                    for i in range(len(texts)):
                        if texts[i]:
                            docs.append(
                                Document(
                                    page_content=texts[i],
                                    metadata={
                                        'name': key,
                                        # 'page':i+1,
                                        'uri': path+parse.quote(key)
                                    }
                                )
                            )
                                    
                elif file_type == 'py' or file_type == 'js':
                    codes = load_code(file_type, key)  # number of functions in the code
                                            
                    if enableParallelSummay=='true':
                        docs = summarize_relevant_codes_using_parallel_processing(codes, key)
                                
                    else:
                        for code in codes:
                            start = code.find('\ndef ')
                            end = code.find(':')                    
                            # print(f'start: {start}, end: {end}')
                                    
                            if start != -1:      
                                function_name = code[start+1:end]
                                # print('function_name: ', function_name)
                                                
                                chat = get_chat(profile_of_LLMs, 0)      
                                        
                                if file_type == 'py':
                                    mode = 'python'
                                elif file_type == 'js':
                                    mode = 'nodejs'
                                else:
                                    mode = file_type  
                                summary = summary_of_code(chat, code, mode)                        
                                            
                                if summary[:len(function_name)]==function_name:
                                    summary = summary[summary.find('\n')+1:len(summary)]
                                                                                        
                                docs.append(
                                    Document(
                                        page_content=summary,
                                        metadata={
                                            'name': key,
                                            # 'page':i+1,
                                            #'uri': path+doc_prefix+parse.quote(key),
                                            'uri': path+key,
                                            'code': code,
                                            'function_name': function_name
                                        }
                                    )
                                )         
                                
                elif file_type == 'png' or file_type == 'jpg' or file_type == 'jpeg':
                    print('extract text from an image: ', key) 
                                       
                    image_obj = s3_client.get_object(Bucket=s3_bucket, Key=key)
                
                    image_content = image_obj['Body'].read()
                    img = Image.open(BytesIO(image_content))
                
                    width, height = img.size 
                    print(f"width: {width}, height: {height}, size: {width*height}")
                
                    isResized = False
                    while(width*height > 5242880):                    
                        width = int(width/2)
                        height = int(height/2)
                        isResized = True
                        print(f"width: {width}, height: {height}, size: {width*height}")
                
                    if isResized:
                        img = img.resize((width, height))
                
                    buffer = BytesIO()
                    img.save(buffer, format="PNG")
                    img_base64 = base64.b64encode(buffer.getvalue()).decode("utf-8")
                                                    
                    # extract text from the image
                    chat = get_chat(profile_of_LLMs, 0)    
                    text = extract_text(chat, img_base64)
                    extracted_text = text[text.find('<result>')+8:len(text)-9] # remove <result> tag
                    print('extracted_text: ', extracted_text)
                    if len(extracted_text)>10:
                        docs.append(
                            Document(
                                page_content=extracted_text,
                                metadata={
                                    'name': key,
                                    # 'page':i+1,
                                    'uri': path+parse.quote(key)
                                }
                            )
                        )        
                                                                                                         
                print('docs size: ', len(docs))
                if len(docs)>0:
                    print('docs[0]: ', docs[0])
                                    
                    ids = store_document_for_opensearch(docs, key)
                    
                create_metadata(bucket=s3_bucket, key=key, meta_prefix=meta_prefix, s3_prefix=s3_prefix, uri=path+parse.quote(key), category=category, documentId=documentId, ids=ids)

            else: # delete if the object is unsupported one for format or size
                try:
                    print('delete the unsupported file: ', key)                                
                    result = s3.delete_object(Bucket=bucket, Key=key)
                    print('result of deletion of the unsupported file: ', result)
                            
                except Exception:
                    err_msg = traceback.format_exc()
                    print('err_msg: ', err_msg)
                    # raise Exception ("Not able to delete unsupported file")
                    
        print('processing time: ', str(time.time() - start_time))
        
        # delete queue
        try:
            sqs.delete_message(QueueUrl=sqsUrl, ReceiptHandle=receiptHandle)
        except Exception as e:        
            print('Fail to delete the queue message: ', e)
            
    return {
        'statusCode': 200
    }
